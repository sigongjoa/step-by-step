import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_correction_pptx(output_path="cognitive_correction_card.pptx"):
    # Create presentation object
    prs = Presentation()
    
    # Set slide size to A4 (8.27 x 11.69 inches)
    prs.slide_width = Inches(8.27)
    prs.slide_height = Inches(11.69)
    
    # Add a blank slide
    slide_layout = prs.slide_layouts[6] # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Define Colors
    INDIGO = RGBColor(79, 70, 229)
    SLATE_900 = RGBColor(15, 23, 42)
    SLATE_400 = RGBColor(148, 163, 184)
    EMERALD = RGBColor(5, 150, 105)
    RED = RGBColor(220, 38, 38)
    WHITE = RGBColor(255, 255, 255)
    
    # 1. Header Section
    # Title Label (Small)
    label_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(0.5), Inches(2.5), Inches(0.3))
    label_box.fill.solid()
    label_box.fill.fore_color.rgb = INDIGO
    label_box.line.visible = False
    
    label_text = label_box.text_frame
    label_text.text = "COGNITIVE CORRECTION REPORT"
    label_text.paragraphs[0].font.size = Pt(8)
    label_text.paragraphs[0].font.bold = True
    label_text.paragraphs[0].font.color.rgb = WHITE
    label_text.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Main Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(5), Inches(0.6))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "인지 교정 분석 카드"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = SLATE_900
    
    # Date/Logo Area
    date_text = slide.shapes.add_textbox(Inches(6.5), Inches(0.9), Inches(1.5), Inches(0.3))
    dp = date_text.text_frame.paragraphs[0]
    dp.text = "2026.01.07"
    dp.font.size = Pt(10)
    dp.font.color.rgb = SLATE_400
    dp.alignment = PP_ALIGN.RIGHT

    # Divider Line
    line = slide.shapes.add_connector(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.6), Inches(7.77), Inches(0.01))
    line.line.color.rgb = SLATE_900
    line.line.width = Pt(1)
    
    # 2. Problem Section
    label_p = slide.shapes.add_textbox(Inches(0.5), Inches(1.9), Inches(4), Inches(0.3))
    lp = label_p.text_frame.paragraphs[0]
    lp.text = "분석 대상 문제 (Target Problem)"
    lp.font.size = Pt(10)
    lp.font.bold = True
    lp.font.color.rgb = SLATE_400
    
    problem_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(2.2), Inches(7.27), Inches(1.5))
    problem_box.fill.solid()
    problem_box.fill.fore_color.rgb = RGBColor(248, 250, 252)
    problem_box.line.color.rgb = RGBColor(241, 245, 249)
    
    problem_text = slide.shapes.add_textbox(Inches(0.7), Inches(2.4), Inches(6.87), Inches(1.1))
    ptf = problem_text.text_frame
    ptf.word_wrap = True
    pp = ptf.paragraphs[0]
    pp.text = "여기에 문제를 입력하세요..."
    pp.font.size = Pt(16)
    pp.font.bold = True
    pp.font.color.rgb = RGBColor(30, 41, 59)
    
    # 3. Diagnostic Section (Bottleneck & Schema)
    # Bottleneck
    label_b = slide.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(3.5), Inches(0.3))
    lb = label_b.text_frame.paragraphs[0]
    lb.text = "취약 지점 (Bottleneck)"
    lb.font.size = Pt(10)
    lb.font.bold = True
    lb.font.color.rgb = INDIGO
    
    b_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.3), Inches(3.5), Inches(1.2))
    b_box.fill.solid()
    b_box.fill.fore_color.rgb = RGBColor(245, 243, 255)
    b_box.line.color.rgb = RGBColor(238, 232, 255)
    
    b_text = slide.shapes.add_textbox(Inches(0.7), Inches(4.5), Inches(3.1), Inches(0.8))
    btf = b_text.text_frame
    bp = btf.paragraphs[0]
    bp.text = "#읽기(Read) #계산(Solve)"
    bp.font.size = Pt(14)
    bp.font.bold = True
    bp.font.color.rgb = INDIGO

    # Schema
    label_s = slide.shapes.add_textbox(Inches(4.27), Inches(4.0), Inches(3.5), Inches(0.3))
    ls = label_s.text_frame.paragraphs[0]
    ls.text = "문제 유형 (Schema)"
    ls.font.size = Pt(10)
    ls.font.bold = True
    ls.font.color.rgb = SLATE_400
    
    s_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.27), Inches(4.3), Inches(3.5), Inches(1.2))
    s_box.fill.solid()
    s_box.fill.fore_color.rgb = RGBColor(248, 250, 252)
    s_box.line.color.rgb = RGBColor(241, 245, 249)
    
    s_text = slide.shapes.add_textbox(Inches(4.47), Inches(4.5), Inches(3.1), Inches(0.8))
    stf = s_text.text_frame
    sp = stf.paragraphs[0]
    sp.text = "#비교하기 #전체와 부분"
    sp.font.size = Pt(14)
    sp.font.bold = True
    sp.font.color.rgb = RGBColor(71, 85, 105)

    # 4. Promises (DO / STOP)
    # DO
    do_icon = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(6.0), Inches(0.8), Inches(0.8))
    do_icon.fill.solid()
    do_icon.fill.fore_color.rgb = EMERALD
    do_icon.line.visible = False
    dit = do_icon.text_frame.paragraphs[0]
    dit.text = "DO"
    dit.font.size = Pt(14)
    dit.font.bold = True
    dit.font.color.rgb = WHITE
    dit.alignment = PP_ALIGN.CENTER
    
    do_desc_box = slide.shapes.add_textbox(Inches(1.5), Inches(6.0), Inches(6.27), Inches(0.8))
    dtf = do_desc_box.text_frame
    d_label = dtf.paragraphs[0]
    d_label.text = "성공을 위한 유일한 행동"
    d_label.font.size = Pt(9)
    d_label.font.color.rgb = EMERALD
    d_label.font.bold = True
    
    d_val = dtf.add_paragraph()
    d_val.text = "예: '모두'라는 단어에 별표 치기"
    d_val.font.size = Pt(18)
    d_val.font.bold = True
    d_val.font.color.rgb = SLATE_900

    # STOP
    stop_icon = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(7.2), Inches(0.8), Inches(0.8))
    stop_icon.fill.solid()
    stop_icon.fill.fore_color.rgb = RED
    stop_icon.line.visible = False
    sit = stop_icon.text_frame.paragraphs[0]
    sit.text = "STOP"
    sit.font.size = Pt(14)
    sit.font.bold = True
    sit.font.color.rgb = WHITE
    sit.alignment = PP_ALIGN.CENTER
    
    stop_desc_box = slide.shapes.add_textbox(Inches(1.5), Inches(7.2), Inches(6.27), Inches(0.8))
    sttf = stop_desc_box.text_frame
    s_label = sttf.paragraphs[0]
    s_label.text = "반드시 멈춰야 할 습관"
    s_label.font.size = Pt(9)
    s_label.font.color.rgb = RED
    s_label.font.bold = True
    
    s_val = sttf.add_paragraph()
    s_val.text = "예: 숫자만 보고 바로 더하지 않기"
    s_val.font.size = Pt(18)
    s_val.font.bold = True
    s_val.font.color.rgb = SLATE_900

    # 5. Footer
    footer = slide.shapes.add_textbox(Inches(0.5), Inches(10.5), Inches(7.27), Inches(0.3))
    fp = footer.text_frame.paragraphs[0]
    fp.text = "MATHESIS LAB | COGNITIVE CORRECTION SYSTEM v1.0"
    fp.font.size = Pt(8)
    fp.font.color.rgb = SLATE_400
    fp.alignment = PP_ALIGN.CENTER
    
    # Save the presentation
    prs.save(output_path)
    print(f"PPTX created successfully: {output_path}")

if __name__ == "__main__":
    create_correction_pptx()
